from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import time

# List of terms to query
terms = [
    "哥本哈根学派",
    "相互依赖武器化",
    "大韩帝国",
    "非洲非殖民化过程",
    "全球南方",
]

OLLAMA_API_URL = "http://localhost:11434/api/generate"  # Default Ollama API endpoint
MAX_RETRIES = 3
RETRY_DELAY = 2

def query_ollama(term, retries=MAX_RETRIES):
    prompt = f"""请以精练的要点形式分析国际关系概念"{term}"300字左右，使用自然语言，减少结构化输出"""

    data = {
        "model": "deepseek-r1:1.5b",  # Change this to your preferred Ollama model
        "prompt": prompt,
        "system": "",
        "stream": False
    }

    for attempt in range(retries):
        try:
            response = requests.post(OLLAMA_API_URL, json=data)
            response.raise_for_status()
            result = response.json()
            response_text = result['response']
            # Remove any content between <think> tags
            import re
            cleaned_text = re.sub(r'<think>.*?</think>', '', response_text, flags=re.DOTALL)
            return cleaned_text.strip()
        except Exception as e:
            if attempt == retries - 1:
                print(f"查询失败 {term}: {str(e)}")
                return f"无法获取'{term}'的解释。错误: {str(e)}"
            time.sleep(RETRY_DELAY)

    return f"在{retries}次尝试后仍无法获取'{term}'的解释"


def format_text(text):
    """格式化文本，确保统一的格式"""
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('•'):
            line = '• ' + line[1:].strip()
        lines.append(line)
    return '\n'.join(lines)


def add_slide(prs, term, explanation):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = term
    title_para.font.size = Pt(32)
    title_para.font.name = '微软雅黑'

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.8), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Add explanation text
    for line in explanation.split('\n'):
        line = line.strip()
        if not line:
            continue

        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.name = '微软雅黑'


def create_ppt(terms):
    prs = Presentation()

    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 添加封面
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = cover_slide.shapes.title
    title.text = "国际关系重要概念解析"

    # 设置封面标题格式
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.name = '微软雅黑'
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)

    if len(cover_slide.placeholders) > 1:
        subtitle = cover_slide.placeholders[1]
        subtitle.text = f"共计{len(terms)}个概念"
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(32)
        subtitle_paragraph.font.name = '微软雅黑'

    # 为每个术语创建幻灯片
    for i, term in enumerate(terms, 1):
        print(f"处理第 {i}/{len(terms)} 个概念: {term}")
        explanation = query_ollama(term)
        print(explanation)
        add_slide(prs, term, explanation)
        print(f"已完成: {term}")

    # 保存文件
    file_path = "国际关系概念解析.pptx"
    prs.save(file_path)
    print(f"\nPPT文件已保存为: {file_path}")
    return file_path

def main():
    try:
        print("开始生成PPT...")
        ppt_file = create_ppt(terms)
        print(f"PPT生成成功！文件位置: {ppt_file}")
    except Exception as e:
        print(f"生成PPT时发生错误: {str(e)}")

if __name__ == "__main__":
    main()