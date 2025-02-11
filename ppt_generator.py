import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import httpx
from openai import OpenAI
import time

# List of terms to query
terms = [
    "哥本哈根学派",
    "相互依赖武器化",
    "大韩帝国",
    "非洲非殖民化过程",
    "全球南方",
    # "第二次结合",
    # "默示谈判和显示谈判",
    # "19世纪70年代清政府处理边疆危机",
    # "欧盟北约共同宣言",
    # "次级制裁"
]

OPENAI_API_KEY = 'sk-l3tKgucSLpTjADcl51C48460D6Eb4eC2B60c2147B25f0b5d'
MAX_RETRIES = 3
RETRY_DELAY = 2


def create_openai_client():
    return OpenAI(
        base_url="https://api.xty.app/v1",
        api_key=OPENAI_API_KEY,
        http_client=httpx.Client(
            base_url="https://api.xty.app/v1",
            follow_redirects=True,
            timeout=30.0
        ),
    )

def query_gpt(client, term, retries=MAX_RETRIES):
    prompt = f"""请以精练的要点形式分析国际关系概念"{term}"，按以下结构分析（用"▪"作为要点符号），300字以内，根据概念灵活增删结构要点，并按照你觉得合理的顺序排列，可能的要点有：
▪ 学术定义
▪ 原因
▪ 背景
▪ 核心主张
▪ 影响
▪ 意义
▪ 重要议题与成果
请确保分析符合国际关系学科特点，请务必删除不必要的要点，注重：
1. 突出权力、利益、制度等国关核心概念
2. 体现国际政治的互动性和复杂性
3. 关注理论与现实的结合
4. 保持学理性和客观性"""

    for attempt in range(retries):
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": """你是一位资深的国际关系学者，精通国际关系理论流派，
                    熟悉现实主义、自由主义、建构主义等主要理论范式，
                    能够从理论和实践层面深入分析国际关系概念。
                    请用专业、准确、简洁的学术语言回答。"""},
                    {"role": "user", "content": prompt}
                ]
            )
            print(response.choices[0].message.content)
            return response.choices[0].message.content
        except Exception as e:
            if attempt == retries - 1:
                print(f"查询失败 {term}: {str(e)}")
                return f"无法获取'{term}'的解释。错误: {str(e)}"
            time.sleep(RETRY_DELAY)

    return f"在{retries}次尝试后仍无法获取'{term}'的解释"



def format_text(text):
    """格式化文本，确保统一的格式"""
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('•'):
            # 确保项目符号格式统一
            line = '• ' + line[1:].strip()
        lines.append(line)
    return '\n'.join(lines)


def add_slide(prs, term, explanation):
    # 使用空白布局
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 设置背景色
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 249, 250)

    # 添加标题 - 调整位置和大小
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))  # 减小标题区域高度
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = term
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.name = '微软雅黑'

    # 添加内容区域 - 上移起始位置
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.8), Inches(4.5))  # 增加内容区域高度
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # 添加内容
    current_section = None
    first_section = True

    for line in explanation.split('\n'):
        line = line.strip()
        if not line:
            continue

        if '【' in line and '】' in line:
            if not first_section:
                spacer = text_frame.add_paragraph()
                spacer.space_before = Pt(1)  # 减小段落间距
            else:
                first_section = False

            p = text_frame.add_paragraph()
            p.text = line.strip('【】')
            p.font.name = '微软雅黑'
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.space_before = Pt(4)  # 减小段落前间距
            p.space_after = Pt(2)  # 减小段落后间距
            current_section = line

        elif line.startswith('▪'):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.name = '微软雅黑'
            p.font.size = Pt(16)
            p.space_before = Pt(1)  # 减小要点间距
            p.space_after = Pt(1)
            p.level = 1


def create_ppt_with_gpt(terms):
    prs = Presentation()
    client = create_openai_client()

    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 添加封面
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = cover_slide.shapes.title
    title.text = "国际关系重要概念解析"

    # 设置封面标题格式
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.name = '微软雅黑'
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)

    if len(cover_slide.placeholders) > 1:
        subtitle = cover_slide.placeholders[1]
        subtitle.text = f"共计{len(terms)}个概念"
        # 设置副标题格式
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(32)
        subtitle_paragraph.font.name = '微软雅黑'

    # 为每个术语创建幻灯片
    for i, term in enumerate(terms, 1):
        print(f"处理第 {i}/{len(terms)} 个概念: {term}")
        explanation = query_gpt(client, term)
        add_slide(prs, term, explanation)
        print(f"已完成: {term}")

    # 保存文件
    file_path = "国际关系概念解析.pptx"
    prs.save(file_path)
    print(f"\nPPT文件已保存为: {file_path}")
    return file_path


def main():
    try:
        print("开始生成PPT...")
        ppt_file = create_ppt_with_gpt(terms)
        print(f"PPT生成成功！文件位置: {ppt_file}")
    except Exception as e:
        print(f"生成PPT时发生错误: {str(e)}")


if __name__ == "__main__":
    main()